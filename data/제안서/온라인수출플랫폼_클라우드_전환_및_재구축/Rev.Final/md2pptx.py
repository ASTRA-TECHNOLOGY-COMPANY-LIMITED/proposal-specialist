#!/usr/bin/env python3
"""
md2pptx.py - Convert proposal markdown to PPTX
Pipeline: Markdown -> HTML -> Chrome Screenshot -> PPTX

Usage:
    python3 md2pptx.py 02_추진전략_및_방법.md
"""

import os, sys, re, subprocess, tempfile, zipfile, math
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Configuration ---
BASE_DIR = Path(__file__).parent
IMAGES_DIR = BASE_DIR / "images"

SLIDE_W, SLIDE_H = 7560000, 10692000
CONTENT_LEFT = 540000
CONTENT_TOP = 1750000
CONTENT_WIDTH = 6480000
CONTENT_BOTTOM = 10100000
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP  # ~8,350,000 EMU

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
HTML_BODY_WIDTH = 700
SCALE_FACTOR = 2

BLUE = RGBColor(0x2E, 0x75, 0xB6)
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x5E)
GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PROJECT_NAME = "온라인수출플랫폼 클라우드 전환 및 재구축"
SECTION_LABEL = "II. 추진전략 및 방법"


# =====================================================
# 1. MARKDOWN PARSER
# =====================================================

def parse_markdown(md_text):
    """Parse markdown into sections by ### headers."""
    sections = []
    current = None
    for line in md_text.split('\n'):
        # Skip top-level # and ## headers (they are section labels, not content)
        if re.match(r'^#{1,2}\s', line) and not line.startswith('###'):
            continue
        if line.startswith('### '):
            if current:
                sections.append(current)
            m = re.match(r'### (\d+\.\d+)\s+(.*)', line)
            badge = m.group(1) if m else ""
            title = m.group(2) if m else line[4:].strip()
            current = {'badge': badge, 'title': title, 'lines': []}
        elif current is not None:
            current['lines'].append(line)
    if current:
        sections.append(current)
    return sections


# =====================================================
# 2. MARKDOWN -> HTML CONVERTER
# =====================================================

def md_inline(text):
    """Bold, italic, code."""
    text = re.sub(r'\*\*([^*]+)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<em>\1</em>', text)
    text = re.sub(r'`([^`]+)`', r'<code>\1</code>', text)
    return text


def table_to_html(rows_raw):
    """Convert list of markdown table lines to HTML table."""
    rows = []
    for line in rows_raw:
        cells = [c.strip() for c in line.strip().strip('|').split('|')]
        rows.append(cells)
    if len(rows) < 3:
        return ''
    header = rows[0]
    data = rows[2:]  # skip separator
    h = '<table>\n<thead><tr>' + ''.join(f'<th>{md_inline(c)}</th>' for c in header) + '</tr></thead>\n<tbody>\n'
    for r in data:
        h += '<tr>' + ''.join(f'<td>{md_inline(c)}</td>' for c in r) + '</tr>\n'
    h += '</tbody></table>'
    return h


def lines_to_html(lines):
    """Convert content lines to HTML body."""
    parts = []
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]
        stripped = line.strip()

        # blank / hr
        if not stripped or stripped == '---':
            i += 1
            continue

        # ##### sub-sub-sub header
        if line.startswith('##### '):
            parts.append(f'<h5>{md_inline(line[6:].strip())}</h5>')
            i += 1
            continue

        # #### sub-sub header
        if line.startswith('#### '):
            parts.append(f'<h4>{md_inline(line[5:].strip())}</h4>')
            i += 1
            continue

        # image ![alt](path)
        img_m = re.match(r'^!\[([^\]]*)\]\(([^)]+)\)', stripped)
        if img_m:
            alt, src = img_m.group(1), img_m.group(2)
            if src.startswith('./images/'):
                src = str(IMAGES_DIR / src[len('./images/'):])
            elif src.startswith('./'):
                src = str(BASE_DIR / src[2:])
            elif not os.path.isabs(src):
                src = str(IMAGES_DIR / src)
            parts.append(f'<div class="img-box"><img src="file://{src}" alt="{alt}"></div>')
            i += 1
            # caption *[그림 ...]*
            if i < n and lines[i].strip().startswith('*['):
                cap = lines[i].strip().strip('*')
                parts.append(f'<p class="caption">{cap}</p>')
                i += 1
            continue

        # table (starts with |, next line has ---)
        if '|' in stripped and stripped.startswith('|'):
            tbl = []
            while i < n and '|' in lines[i].strip() and lines[i].strip().startswith('|'):
                tbl.append(lines[i])
                i += 1
            parts.append(table_to_html(tbl))
            continue

        # blockquote >
        if stripped.startswith('>'):
            bq = []
            while i < n and lines[i].strip().startswith('>'):
                bq.append(md_inline(lines[i].strip().lstrip('>').strip()))
                i += 1
            parts.append('<blockquote>' + '<br>'.join(bq) + '</blockquote>')
            continue

        # unordered list - or *
        if re.match(r'^[-*]\s', stripped):
            items = []
            while i < n and re.match(r'^\s*[-*]\s', lines[i]):
                items.append(md_inline(re.sub(r'^\s*[-*]\s+', '', lines[i])))
                i += 1
            parts.append('<ul>' + ''.join(f'<li>{it}</li>' for it in items) + '</ul>')
            continue

        # numbered list
        if re.match(r'^\d+\.\s', stripped):
            items = []
            while i < n and re.match(r'^\s*\d+\.\s', lines[i]):
                items.append(md_inline(re.sub(r'^\s*\d+\.\s+', '', lines[i])))
                i += 1
            parts.append('<ol>' + ''.join(f'<li>{it}</li>' for it in items) + '</ol>')
            continue

        # paragraph (collect consecutive non-special lines)
        para = []
        while i < n:
            l = lines[i]
            s = l.strip()
            if not s or s == '---' or s.startswith('#') or s.startswith('|') or s.startswith('>') or re.match(r'^!\[', s) or re.match(r'^[-*]\s', s) or re.match(r'^\d+\.\s', s):
                break
            para.append(s)
            i += 1
        if para:
            parts.append(f'<p>{md_inline(" ".join(para))}</p>')

    return '\n'.join(parts)


def wrap_html(body):
    return f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
body {{
  font-family: 'Malgun Gothic','Apple SD Gothic Neo','Noto Sans KR',sans-serif;
  width: {HTML_BODY_WIDTH}px; margin:0; padding:20px;
  font-size:11pt; line-height:1.65; color:#333; background:#fff;
}}
h4 {{
  font-size:12.5pt; color:#1A3C5E; margin:22px 0 8px;
  border-left:4px solid #2E75B6; padding-left:10px; font-weight:700;
}}
h5 {{
  font-size:11.5pt; color:#1A3C5E; margin:18px 0 6px; font-weight:700;
}}
p {{ margin:6px 0; text-align:justify; }}
table {{
  width:100%; border-collapse:collapse; margin:10px 0; font-size:9.5pt;
}}
th {{
  background:#2E75B6; color:#fff; padding:8px 10px; text-align:center;
  font-weight:700; border:1px solid #2066A7;
}}
td {{
  padding:7px 10px; border:1px solid #D0D0D0; vertical-align:top;
}}
tbody tr:nth-child(even) {{ background:#F5F8FC; }}
strong {{ color:#1A3C5E; }}
blockquote {{
  background:#F0F5FA; border-left:4px solid #2E75B6;
  padding:12px 16px; margin:12px 0; font-size:10pt; color:#444;
}}
.img-box {{ text-align:center; margin:14px 0 4px; }}
.img-box img {{ max-width:100%; height:auto; }}
.caption {{ text-align:center; font-size:9pt; color:#666; margin:2px 0 14px; }}
ul,ol {{ margin:6px 0; padding-left:24px; }}
li {{ margin:3px 0; }}
code {{ background:#F0F0F0; padding:1px 4px; border-radius:3px; font-size:9pt; }}
</style></head>
<body>{body}</body></html>'''


# =====================================================
# 3. CHROME HEADLESS RENDERER
# =====================================================

def render_to_png(html_str, png_path, temp_dir):
    html_file = os.path.join(temp_dir, 'render.html')
    with open(html_file, 'w', encoding='utf-8') as f:
        f.write(html_str)
    win_w = HTML_BODY_WIDTH + 40
    cmd = [
        CHROME, '--headless=new', '--disable-gpu', '--no-sandbox',
        f'--window-size={win_w},4000',
        f'--force-device-scale-factor={SCALE_FACTOR}',
        f'--screenshot={png_path}',
        '--default-background-color=FFFFFFFF',
        f'file://{html_file}'
    ]
    subprocess.run(cmd, capture_output=True, timeout=60)
    # Trim whitespace at bottom
    if os.path.exists(png_path):
        trim_bottom(png_path)
    return png_path


def trim_bottom(png_path):
    """Remove trailing white rows from PNG."""
    img = Image.open(png_path)
    px = img.load()
    w, h = img.size
    bottom = h - 1
    while bottom > 0:
        row_white = all(
            px[x, bottom][:3] == (255, 255, 255) or px[x, bottom][:3] == (0, 0, 0) and px[x, bottom][3] == 0
            for x in range(0, w, max(1, w // 20))
        )
        if not row_white:
            break
        bottom -= 1
    if bottom < h - 5:
        img.crop((0, 0, w, bottom + 5)).save(png_path)
    img.close()


# =====================================================
# 4. PPTX BUILDER
# =====================================================

def add_chrome(slide, badge, title, page_num):
    """Add slide chrome: stripe, triangle, header, badge, title, divider, page."""

    # Left blue stripe
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 180000, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

    # Top-left triangle (approximate with right triangle)
    t = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, 720000, 540000)
    t.fill.solid(); t.fill.fore_color.rgb = BLUE; t.line.fill.background()
    # Flip vertically to point down-right
    t.rotation = 180.0

    # Project name (top right)
    tb = slide.shapes.add_textbox(2520000, 288000, 4680000, 216000)
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = PROJECT_NAME; r.font.size = Pt(8); r.font.color.rgb = GRAY

    # Section label (below project name)
    tb = slide.shapes.add_textbox(2520000, 468000, 4680000, 216000)
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = SECTION_LABEL; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = DARK_BLUE

    # Page number background (red rounded rect)
    pg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 7128000, 10260000, 251999, 251999)
    pg.fill.solid(); pg.fill.fore_color.rgb = RED; pg.line.fill.background()
    # Page number text
    tb = slide.shapes.add_textbox(7128000, 10260000, 251999, 251999)
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(page_num); r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = WHITE

    # Badge background (blue rounded rect)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 540000, 1260000, 432000, 324000)
    bg.fill.solid(); bg.fill.fore_color.rgb = BLUE; bg.line.fill.background()
    # Badge text
    tb = slide.shapes.add_textbox(540000, 1260000, 432000, 324000)
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = badge
    r.font.size = Pt(11) if len(badge) > 2 else Pt(16)
    r.font.bold = True; r.font.color.rgb = WHITE

    # Title
    tb = slide.shapes.add_textbox(1080000, 1260000, 5400000, 324000)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DARK_BLUE

    # Divider
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 540000, 1620000, 6480000, 18000)
    d.fill.solid(); d.fill.fore_color.rgb = BLUE; d.line.fill.background()


def build_pptx(sections, temp_dir, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    page_num = 1

    for sec in sections:
        png = sec.get('png')
        if not png or not os.path.exists(png):
            continue

        img = Image.open(png)
        iw, ih = img.size
        scale = CONTENT_WIDTH / iw
        ih_emu = int(ih * scale)

        if ih_emu <= CONTENT_HEIGHT:
            sl = prs.slides.add_slide(blank)
            add_chrome(sl, sec['badge'], sec['title'], page_num)
            sl.shapes.add_picture(png, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, ih_emu)
            page_num += 1
        else:
            px_per_slide = int(CONTENT_HEIGHT / scale)
            num = math.ceil(ih / px_per_slide)
            for s in range(num):
                y0 = s * px_per_slide
                y1 = min((s + 1) * px_per_slide, ih)
                crop = img.crop((0, y0, iw, y1))
                cp = os.path.join(temp_dir, f'crop_{sec["badge"]}_{s}.png')
                crop.save(cp)
                ch = int((y1 - y0) * scale)

                sl = prs.slides.add_slide(blank)
                add_chrome(sl, sec['badge'], sec['title'], page_num)
                sl.shapes.add_picture(cp, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, ch)
                page_num += 1

        img.close()

    # Save with ZIP cleanup
    tmp = os.path.join(temp_dir, '_raw.pptx')
    prs.save(tmp)
    seen = {}
    with zipfile.ZipFile(tmp) as zin:
        for item in zin.infolist():
            seen[item.filename] = item
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name, info in seen.items():
                zout.writestr(info, zin.read(info.filename))

    print(f'\n=== Created: {output_path} ({page_num - 1} slides) ===')


# =====================================================
# 5. MAIN
# =====================================================

def main():
    md_file = sys.argv[1] if len(sys.argv) > 1 else '02_추진전략_및_방법.md'
    md_path = BASE_DIR / md_file
    out_path = str(md_path.with_suffix('.pptx'))

    print(f'Input:  {md_path}')
    print(f'Output: {out_path}')

    md_text = md_path.read_text(encoding='utf-8')
    sections = parse_markdown(md_text)
    print(f'Parsed {len(sections)} sections')

    with tempfile.TemporaryDirectory() as tmp:
        for i, sec in enumerate(sections):
            html_body = lines_to_html(sec['lines'])
            html_full = wrap_html(html_body)
            png = os.path.join(tmp, f's{i:02d}.png')
            render_to_png(html_full, png, tmp)
            sec['png'] = png
            # Get image size for info
            if os.path.exists(png):
                im = Image.open(png)
                print(f'  [{sec["badge"]}] {sec["title"]}  ->  {im.size[0]}x{im.size[1]}px')
                im.close()
            else:
                print(f'  [{sec["badge"]}] {sec["title"]}  ->  RENDER FAILED')

        build_pptx(sections, tmp, out_path)


if __name__ == '__main__':
    main()
